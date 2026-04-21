"""JP-Force 社内テンプレート スライド自動生成ツール

使い方:
    python slide_generator.py <spec.json> [--out output.pptx]
    echo '{"slides":[...]}' | python slide_generator.py - --out output.pptx

スペック(JSON)の仕様は example_spec.json と各ビルダー関数の docstring を参照。
Claude デスクトップアプリから呼び出す際は、本ファイルと同階層に
assets/jpf_logo.png を置いておくと表紙ロゴとして自動挿入される。
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- テーマ定義（社内テンプレ準拠） ----------
FONT_JP = "Noto Sans JP"

class C:
    GOLD = RGBColor(0xA9, 0x96, 0x5B)
    GOLD_LIGHT = RGBColor(0xD4, 0xB9, 0x7A)
    GOLD_BG = RGBColor(0xED, 0xE8, 0xDF)
    NAVY = RGBColor(0x00, 0x1B, 0x3A)
    DARK = RGBColor(0x43, 0x43, 0x43)
    MID = RGBColor(0x55, 0x5F, 0x6B)
    MID2 = RGBColor(0x3D, 0x4B, 0x5F)
    GRAY = RGBColor(0xC2, 0xCC, 0xD2)
    BG = RGBColor(0xF5, 0xF5, 0xF5)
    SOFT_BG = RGBColor(0xFA, 0xF9, 0xF5)  # 淡い白（わずかに温かい）
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN = Inches(0.4)
ASSETS = Path(__file__).parent / "assets"


# ---------- ユーティリティ ----------
def _set_text(tf, text: str, size: float, *, bold=False, color=C.DARK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_JP):
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(slide, left, top, width, height, text, size, **kw):
    box = slide.shapes.add_textbox(left, top, width, height)
    _set_text(box.text_frame, text, size, **kw)
    return box


def _add_rect(slide, left, top, width, height, fill, *, line=False, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = fill
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _set_flip_rot(shape, *, flip_h=False, flip_v=False, rotation=None):
    """add_shape では設定できない flip / rotation を XML 直接設定"""
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    if flip_h:
        xfrm.set('flipH', '1')
    if flip_v:
        xfrm.set('flipV', '1')
    if rotation is not None:
        xfrm.set('rot', str(int(rotation * 60000)))  # 度 → 60000分の1度


def _add_footer(slide, page_no: int, footer_title: str):
    # 下部ゴールドライン（スライド幅いっぱい）
    _add_rect(slide, Emu(0), Inches(5.45), SLIDE_W, Inches(0.03), C.GOLD)
    _add_text(slide, Inches(0.30), Inches(5.49), Inches(4), Inches(0.13),
              "RESTRICTED / CONFIDENTIAL", 6.5, color=C.GRAY)
    _add_text(slide, Inches(3.50), Inches(5.49), Inches(3.5), Inches(0.13),
              footer_title, 6.5, color=C.GRAY, align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(9.20), Inches(5.49), Inches(0.6), Inches(0.13),
              str(page_no), 7, color=C.GRAY, bold=True, align=PP_ALIGN.RIGHT)


def _page_header(slide, title: str, subtitle: str | None = None):
    # ゴールド縦バー + タイトル
    _add_rect(slide, MARGIN, Inches(0.42), Inches(0.08), Inches(0.35), C.GOLD)
    _add_text(slide, Inches(0.58), Inches(0.4), Inches(9), Inches(0.45),
              title, 18, bold=True, color=C.NAVY)
    if subtitle:
        _add_text(slide, Inches(0.58), Inches(0.82), Inches(9), Inches(0.3),
                  subtitle, 9, color=C.MID)
    # 区切り線
    _add_rect(slide, MARGIN, Inches(1.18), Inches(9.2), Inches(0.008), C.GRAY)


# ---------- 各スライド ビルダー ----------
def build_cover(slide, spec: dict):
    eyebrow = spec.get("eyebrow", "Market Research Report")
    title = spec.get("title", "タイトル")
    subtitle = spec.get("subtitle", "")
    date = spec.get("date", "")
    show_logo = spec.get("logo", True)

    # 背景
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, C.SOFT_BG)

    # --- 装飾シェイプ（参考テンプレ準拠） ---
    # グレーのスラント型パネル（右側）: flowChartManualInput を -90° + flipH
    gray = _add_rect(slide, Inches(4.808), Inches(0.430),
                     Inches(5.625), Inches(4.765), C.GRAY,
                     shape=MSO_SHAPE.FLOWCHART_MANUAL_INPUT)
    _set_flip_rot(gray, flip_h=True, rotation=-90)
    # ゴールドの右下三角: rtTriangle を flipH
    gold_tri = _add_rect(slide, Inches(3.681), Inches(2.073),
                         Inches(6.319), Inches(3.552), C.GOLD,
                         shape=MSO_SHAPE.RIGHT_TRIANGLE)
    _set_flip_rot(gold_tri, flip_h=True)

    # 左端基準 (参考テンプレ準拠: 0.25 inch)
    left = Inches(0.25)
    # Eyebrow：ゴールド背景バー + ネイビー文字（縦中央配置: ブロック上端 T=1.60）
    _add_rect(slide, left, Inches(1.60), Inches(2.3), Inches(0.3), C.GOLD)
    _add_text(slide, left, Inches(1.60), Inches(2.3), Inches(0.3),
              eyebrow, 8.5, bold=True, color=C.NAVY,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # メインタイトル（縦中心 ≒ 2.81 inch に配置）
    _add_text(slide, left, Inches(2.07), Inches(6), Inches(0.9),
              title, 38, bold=True, color=C.DARK)
    # サブタイトル
    _add_text(slide, left, Inches(2.93), Inches(6), Inches(0.6),
              subtitle, 24, color=C.GOLD_LIGHT)
    # 区切りゴールドライン
    _add_rect(slide, left, Inches(3.65), Inches(3.5), Inches(0.03), C.GOLD)
    # 日付
    _add_text(slide, left, Inches(3.77), Inches(4), Inches(0.25),
              date, 10, color=C.DARK)
    # ロゴ（任意・右下）
    logo = ASSETS / "jpf_logo.png"
    if show_logo and logo.exists():
        slide.shapes.add_picture(str(logo), Inches(8.6), Inches(5.05),
                                 height=Inches(0.3))


def build_toc(slide, spec: dict, page_no: int, footer_title: str):
    items = spec.get("items", [])
    heading = spec.get("heading", "目次")

    # 背景（淡い白）
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, C.SOFT_BG)
    # ヘッダ: 左端ゴールドバー + 見出し
    _add_rect(slide, Emu(0), Emu(0), Inches(0.08), Inches(0.60), C.GOLD)
    _add_text(slide, Inches(0.22), Emu(0), Inches(2), Inches(0.60),
              heading, 15, bold=True, color=C.DARK, anchor=MSO_ANCHOR.MIDDLE)

    # 2列 × 最大5行レイアウト
    card_w = Inches(4.65)
    card_h = Inches(0.56)
    row_gap = Inches(0.08)  # 0.64 - 0.56
    start_y = Inches(1.25)
    col_x = [Inches(0.28), Inches(5.08)]

    for idx, it in enumerate(items):
        col = idx // 5
        row = idx % 5
        if col > 1:
            break
        x = col_x[col]
        y = start_y + row * (card_h + row_gap)
        # 白カード背景
        _add_rect(slide, x, y, card_w, card_h, C.WHITE)
        # 左ゴールドアクセントバー
        _add_rect(slide, x, y, Inches(0.06), card_h, C.GOLD)
        # 番号（ゴールド太字）
        _add_text(slide, x + Inches(0.14), y, Inches(0.55), card_h,
                  it.get("no", f"{idx+1:02d}"), 18, bold=True, color=C.GOLD,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル（ネイビー太字）
        tx = x + Inches(0.68)
        tw = Inches(3.89)
        _add_text(slide, tx, y + Inches(0.05), tw, Inches(0.26),
                  it.get("title", ""), 10.5, bold=True, color=C.NAVY)
        # 説明（ミッドグレー）
        _add_text(slide, tx, y + Inches(0.30), tw, Inches(0.22),
                  it.get("desc", ""), 8, color=C.MID)

    _add_footer(slide, page_no, footer_title)


def build_section(slide, spec: dict, page_no: int, footer_title: str):
    number = spec.get("number", "01")
    title = spec.get("title", "Section Title")
    subtitle = spec.get("subtitle", "")

    # 全面ゴールド背景
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, C.GOLD)
    # 左側・白い正方形アクセント
    _add_rect(slide, Inches(0.38), Inches(2.44),
              Inches(0.74), Inches(0.74), C.WHITE)
    # 番号（白四角内・ゴールド太字）
    _add_text(slide, Inches(0.38), Inches(2.44),
              Inches(0.74), Inches(0.74),
              number, 23, bold=True, color=C.GOLD,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # タイトル（白文字）
    _add_text(slide, Inches(1.49), Inches(2.44),
              Inches(7.81), Inches(0.74),
              title, 24, bold=True, color=C.WHITE,
              anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_text(slide, Inches(1.49), Inches(3.25),
                  Inches(7.81), Inches(0.4),
                  subtitle, 11, color=C.WHITE)
    # フッター（白文字）
    _add_text(slide, Inches(0.49), Inches(5.29),
              Inches(2), Inches(0.1),
              "RESTRICTED DISTRIBUTION", 5, bold=True, color=C.WHITE)
    _add_text(slide, Inches(0.30), Inches(5.37),
              Inches(6), Inches(0.13),
              f"© 2026 JP-Force. All rights reserved.  |  {footer_title}",
              5, color=C.WHITE)
    _add_text(slide, Inches(9.20), Inches(5.05),
              Inches(0.60), Inches(0.2),
              str(page_no), 8, bold=True, color=C.WHITE, align=PP_ALIGN.RIGHT)
    # ロゴ（右下・小さめ）
    logo = ASSETS / "jpf_logo.png"
    if spec.get("logo", True) and logo.exists():
        slide.shapes.add_picture(str(logo), Inches(8.70), Inches(5.30),
                                 height=Inches(0.16))


# ---------- コンテンツブロック ----------
def _block_kpi_row(slide, y: Emu, items: list[dict]) -> Emu:
    n = max(1, len(items))
    gap = Inches(0.15)
    total_w = Inches(9.2)
    card_w = Emu((total_w - gap * (n - 1)) // n)
    card_h = Inches(1.25)
    for i, it in enumerate(items):
        x = MARGIN + i * (card_w + gap)
        _add_rect(slide, x, y, card_w, card_h, C.WHITE, line=True)
        _add_rect(slide, x, y, Inches(0.05), card_h, C.GOLD)
        _add_text(slide, x + Inches(0.2), y + Inches(0.15), card_w, Inches(0.55),
                  str(it.get("value", "")), 28, bold=True, color=C.NAVY)
        unit = it.get("unit", "")
        if unit:
            _add_text(slide, x + Inches(0.2), y + Inches(0.68), card_w, Inches(0.25),
                      unit, 10, bold=True, color=C.GOLD)
        _add_text(slide, x + Inches(0.2), y + Inches(0.92), card_w - Inches(0.3), Inches(0.3),
                  it.get("label", ""), 8.5, color=C.MID)
    return y + card_h + Inches(0.2)


def _block_bullets(slide, y: Emu, data: dict) -> Emu:
    heading = data.get("heading")
    items = data.get("items", [])
    if heading:
        _add_text(slide, MARGIN, y, Inches(9.2), Inches(0.3),
                  heading, 11, bold=True, color=C.NAVY)
        y = y + Inches(0.32)
    for it in items:
        _add_rect(slide, MARGIN + Inches(0.05), y + Inches(0.1),
                  Inches(0.08), Inches(0.08), C.GOLD, shape=MSO_SHAPE.OVAL)
        _add_text(slide, MARGIN + Inches(0.25), y, Inches(9), Inches(0.3),
                  str(it), 10, color=C.DARK)
        y = y + Inches(0.3)
    return y + Inches(0.1)


def _block_numbered(slide, y: Emu, data: dict) -> Emu:
    heading = data.get("heading")
    items = data.get("items", [])
    cols = data.get("cols", 2)
    if heading:
        _add_text(slide, MARGIN, y, Inches(9.2), Inches(0.3),
                  heading, 11, bold=True, color=C.NAVY)
        y = y + Inches(0.32)
    gap_x = Inches(0.15)
    total_w = Inches(9.2)
    card_w = Emu((total_w - gap_x * (cols - 1)) // cols)
    card_h = Inches(0.7)
    start_y = y
    rows = (len(items) + cols - 1) // cols
    for i, it in enumerate(items):
        r, c = i // cols, i % cols
        x = MARGIN + c * (card_w + gap_x)
        yy = start_y + r * (card_h + Inches(0.1))
        _add_rect(slide, x, yy, card_w, card_h, C.BG)
        _add_rect(slide, x, yy, Inches(0.6), card_h, C.NAVY)
        _add_text(slide, x, yy, Inches(0.6), card_h,
                  it.get("no", f"{i+1:02d}"), 14, bold=True, color=C.GOLD,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = x + Inches(0.7)
        tw = card_w - Inches(0.75)
        title = it.get("title", "")
        text = it.get("text", "")
        if title:
            _add_text(slide, tx, yy + Inches(0.05), tw, Inches(0.3),
                      title, 10, bold=True, color=C.NAVY)
            _add_text(slide, tx, yy + Inches(0.32), tw, Inches(0.35),
                      text, 8.5, color=C.DARK)
        else:
            _add_text(slide, tx, yy, tw, card_h,
                      text, 9, color=C.DARK, anchor=MSO_ANCHOR.MIDDLE)
    return start_y + rows * (card_h + Inches(0.1)) + Inches(0.1)


def _block_table(slide, y: Emu, data: dict) -> Emu:
    headers = data.get("headers", [])
    rows = data.get("rows", [])
    heading = data.get("heading")
    if heading:
        _add_text(slide, MARGIN, y, Inches(9.2), Inches(0.3),
                  heading, 11, bold=True, color=C.NAVY)
        y = y + Inches(0.32)
    cols = len(headers)
    if cols == 0:
        return y
    total_w = Inches(9.2)
    row_h = Inches(0.35)
    col_w = Emu(total_w // cols)
    # ヘッダ
    _add_rect(slide, MARGIN, y, total_w, row_h, C.NAVY)
    for i, h in enumerate(headers):
        _add_text(slide, MARGIN + i * col_w, y, col_w, row_h,
                  str(h), 9.5, bold=True, color=C.WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y = y + row_h
    for ri, row in enumerate(rows):
        bg = C.WHITE if ri % 2 == 0 else C.BG
        _add_rect(slide, MARGIN, y, total_w, row_h, bg)
        for i, cell in enumerate(row):
            _add_text(slide, MARGIN + i * col_w + Inches(0.1), y, col_w - Inches(0.2), row_h,
                      str(cell), 9, color=C.DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = y + row_h
    # ボーダー
    _add_rect(slide, MARGIN, y, total_w, Inches(0.015), C.GOLD)
    return y + Inches(0.15)


def _block_text(slide, y: Emu, data: dict) -> Emu:
    content = data.get("content", "")
    size = data.get("size", 10)
    color = getattr(C, data.get("color", "DARK"), C.DARK)
    h = Inches(data.get("height", 0.5))
    _add_text(slide, MARGIN, y, Inches(9.2), h,
              content, size, color=color,
              bold=data.get("bold", False))
    return y + h + Inches(0.1)


def _block_callout(slide, y: Emu, data: dict) -> Emu:
    content = data.get("content", "")
    h = Inches(data.get("height", 0.6))
    _add_rect(slide, MARGIN, y, Inches(9.2), h, C.GOLD_BG)
    _add_rect(slide, MARGIN, y, Inches(0.08), h, C.GOLD)
    _add_text(slide, MARGIN + Inches(0.2), y, Inches(9), h,
              content, 10, bold=True, color=C.NAVY, anchor=MSO_ANCHOR.MIDDLE)
    return y + h + Inches(0.15)


BLOCKS: dict[str, Callable] = {
    "kpi_row": lambda s, y, d: _block_kpi_row(s, y, d.get("items", [])),
    "bullets": _block_bullets,
    "numbered": _block_numbered,
    "table": _block_table,
    "text": _block_text,
    "callout": _block_callout,
}


def _estimate_block_height(block: dict) -> Emu:
    """ブロック描画後の縦方向消費量を推定（描画前オーバーフロー検知用）"""
    t = block.get("type")
    if t == "kpi_row":
        return Inches(1.25 + 0.2)
    if t == "bullets":
        h = Inches(0.4) if block.get("heading") else Emu(0)
        h += Inches(0.3) * len(block.get("items", []))
        return h + Inches(0.1)
    if t == "numbered":
        cols = block.get("cols", 2)
        n = len(block.get("items", []))
        rows = (n + cols - 1) // cols
        h = Inches(0.4) if block.get("heading") else Emu(0)
        h += rows * Inches(0.8)
        return h + Inches(0.1)
    if t == "table":
        h = Inches(0.4) if block.get("heading") else Emu(0)
        n = 1 + len(block.get("rows", []))  # ヘッダ + 行
        return h + n * Inches(0.35) + Inches(0.15)
    if t == "callout":
        return Inches(block.get("height", 0.6)) + Inches(0.15)
    if t == "text":
        return Inches(block.get("height", 0.5)) + Inches(0.1)
    return Inches(0.5)


def build_content(slide, spec: dict, page_no: int, footer_title: str):
    title = spec.get("title", "")
    subtitle = spec.get("subtitle", "")
    # 背景（淡い白）
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, C.SOFT_BG)
    _page_header(slide, title, subtitle)

    safe_bottom = Inches(5.30)  # フッターライン 5.45 から 0.15 inch 余白
    y = Inches(1.35)
    skipped: list[str] = []
    for block in spec.get("blocks", []):
        btype = block.get("type", "?")
        fn = BLOCKS.get(btype)
        est = _estimate_block_height(block)
        # 事前オーバーフロー検知: フッター領域を侵食する場合はスキップ
        if y + est > safe_bottom:
            skipped.append(btype)
            continue
        if fn is None:
            y = _block_text(slide, y, {"content": f"[未対応ブロック: {btype}]"})
        else:
            y = fn(slide, y, block)

    if skipped:
        print(
            f"[WARN] slide {page_no} '{title}': "
            f"領域不足により {len(skipped)} ブロック省略 ({', '.join(skipped)}) "
            f"— spec を別スライドに分割してください",
            file=sys.stderr,
        )
    _add_footer(slide, page_no, footer_title)


# ---------- メイン ----------
SLIDE_BUILDERS: dict[str, Callable] = {
    "cover": build_cover,
    "toc": build_toc,
    "section": build_section,
    "content": build_content,
}


def generate(spec: dict, out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank

    footer_title = spec.get("title_short", "JP-Force Report")
    for i, s in enumerate(spec.get("slides", []), start=1):
        slide = prs.slides.add_slide(blank)
        builder = SLIDE_BUILDERS.get(s.get("type", "content"))
        if builder is None:
            continue
        if s["type"] == "cover":
            builder(slide, s)
        else:
            builder(slide, s, i, footer_title)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main():
    ap = argparse.ArgumentParser(description="JP-Force 社内テンプレ スライド自動生成")
    ap.add_argument("spec", help="スペックJSONファイル（- で標準入力）")
    ap.add_argument("--out", default=None,
                    help="出力先 .pptx パス（未指定なら output/<spec名>.pptx）")
    args = ap.parse_args()

    if args.spec == "-":
        spec = json.load(sys.stdin)
        default_name = "slides.pptx"
    else:
        spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
        default_name = Path(args.spec).stem + ".pptx"

    out_dir = Path(__file__).parent / "output"
    out_path = Path(args.out) if args.out else out_dir / default_name
    if not out_path.is_absolute() and args.out is None:
        out_path = out_dir / default_name

    out = generate(spec, out_path)
    print(f"[OK] 生成完了: {out}")


if __name__ == "__main__":
    main()
